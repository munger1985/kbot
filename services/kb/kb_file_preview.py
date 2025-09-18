import os
import mimetypes
import base64
from PyPDF2 import PdfReader, PdfWriter
import docx
import openpyxl
from pptx import Presentation
from io import BytesIO
from typing import Any
from dao.repositories.kbot_md_kb_files_repo import KbotMdKbFilesRepository
import tempfile



# 文件预览处理器
class FilePreview:

    async def get_preview(self, 
                    file_id: str, 
                    max_text_length: int = 10000, 
                    max_pages: int = 10, 
                    max_sheets: int = 5, 
                    max_slides: int = 10, 
                    pdf_pages: int | list[int] | None = None, 
                    word_page: int | None = None,
                    sheet_index: int = 0, 
                    start_index: int = 0,
                    slide: int | None = None, 
                    ) -> dict[str, Any]:
        """
        获取文件预览数据
        
        Args:
            file_id: 文件ID
            max_text_length: 最大文本长度
            max_pages: 最大PDF页数
            max_sheets: 最大Excel工作表数
            max_slides: 最大PPT幻灯片数
            page: 指定Word/PDF页码
            sheet: 指定Excel工作表名
            slide: 指定PPT幻灯片编号
            
        Returns:
            包含预览数据的字典
        """
        # 从数据库获取文件路径
        file_path = await KbotMdKbFilesRepository().get_path_by_id(file_id)

        if not file_path:
            return {"error": "文件不存在", "file_id": file_id}
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            return {"error": "文件不存在", "file_path": file_path}
        
        # 获取文件MIME类型
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            mime_type = 'application/octet-stream'
        
        # 根据文件类型调用不同的预览方法
        try:
            if mime_type == 'text/plain': 
                preview_data = await self._preview_text(file_path, max_text_length)
            elif mime_type in ['image/jpeg', 'image/png', 'image/gif']:
                preview_data = await self._preview_image(file_path)
            elif mime_type == 'application/pdf':
                preview_data = await self._preview_pdf(file_path, pdf_pages)
            elif mime_type in ['application/vnd.ms-excel', 
                              'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
                preview_data = await self._preview_excel_document(file_path, max_sheets, start_index, sheet_index)
            elif mime_type in ['application/vnd.ms-powerpoint',
                              'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                preview_data = await self._preview_ppt_document(file_path, max_slides, slide)
            elif mime_type in ['application/msword', 
                              'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
                preview_data = await self._preview_word_document(file_path, max_pages, word_page)
            elif mime_type == 'video/': 
                preview_data = await self._preview_video(file_path)
            elif mime_type == 'audio/': 
                preview_data = await self._preview_audio(file_path)
            else:
                pass
                
            preview_data.update({
                "file_id": file_id,
                "file_name": os.path.basename(file_path),
                "mime_type": mime_type,
                "file_size": os.path.getsize(file_path),
                "success": True
            })
            return preview_data
        except Exception as e:
            return {
                "error": f"处理文件时出错: {str(e)}",
                "file_id": file_id,
                "file_name": os.path.basename(file_path),
                "success": False
            }
    
    async def _preview_text(self, file_path: str, max_length: int) -> dict[str, Any]:
        """预览文本文件"""
        file_size = os.path.getsize(file_path)
        
        # 读取文件内容
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read(max_length)
        except UnicodeDecodeError:
            # 如果UTF-8解码失败，尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    content = f.read(max_length)
            except:
                content = "无法解码文本内容"
        
        truncated = file_size > max_length
        if truncated:
            content += "\n...(文件过大，只显示部分内容)"
        
        return {
            "preview_type": "text",
            "content": content
        }
    
    async def _preview_image(self, file_path: str) -> dict[str, Any]:
        """预览图片文件"""
        # 将图片转换为base64编码
        with open(file_path, 'rb') as f:
            image_data = base64.b64encode(f.read()).decode('utf-8')
        
        return {
            "preview_type": "image",
            "content": image_data
        }
    
    async def _preview_pdf(self, file_path: str, pages: int | list[int] | None = None) -> dict[str, Any]:
        """预览PDF文件，提取指定页的二进制内容"""
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PdfReader(f)
                total_pages = len(pdf_reader.pages)
                
                # 处理页码参数
                if pages is None:
                    target_pages = list(range(1, total_pages + 1))
                elif isinstance(pages, int):
                    target_pages = [pages]
                else:
                    target_pages = pages
                
                # 验证页码有效性
                valid_pages = []
                for page_num in target_pages:
                    if 1 <= page_num <= total_pages:
                        valid_pages.append(page_num)
                
                if not valid_pages:
                    return {
                        "preview_type": "pdf",
                        "error": f"无效的页码: {pages}. 有效范围: 1-{total_pages}"
                    }
                
                # 提取指定页的内容
                pages_content = []
                
                for page_num in valid_pages:
                    page_index = page_num - 1
                    
                    # 方法1: 提取单页为新的PDF二进制数据
                    writer = PdfWriter()
                    writer.add_page(pdf_reader.pages[page_index])
                    
                    # 将单页PDF写入内存缓冲区  
                    buffer = BytesIO()
                    writer.write(buffer)
                    pdf_binary = buffer.getvalue()
                    
                    # 转换为Base64编码以便传输
                    pdf_base64 = base64.b64encode(pdf_binary).decode('utf-8')
                    
                    pages_content.append({
                        "page_number": page_num,
                        "pdf_binary": pdf_base64,  # 单页PDF的Base64编码
                        "binary_size": len(pdf_binary)  # 二进制数据大小
                    })
                
                # 返回结果
                return {
                    "preview_type": "pdf",
                    "total_pages": total_pages,
                    "content": pages_content,
                    "message": f"成功提取 {len(valid_pages)} 页内容"
                }
                    
        except Exception as e:
            return {
                "preview_type": "pdf",
                "error": f"处理PDF文件时出错: {str(e)}"
            }
    
    async def _preview_word_document(self, file_path: str, max_pages: int, page: int | None) -> dict[str, Any]:
        """预览Word文档 - 按页提取内容"""
        try:
            doc = docx.Document(file_path)
            
            # 获取文档的所有段落
            all_paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():  # 只处理非空段落
                    all_paragraphs.append(para.text)
            
            # 获取文档的所有表格内容
            all_tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                all_tables.append(table_data)
            
            # 模拟分页 - 每10个段落为一页
            pages = []
            paragraphs_per_page = 10
            total_pages = (len(all_paragraphs) + paragraphs_per_page - 1) // paragraphs_per_page
            
            # 如果指定了页码，只提取该页
            if page is not None and 1 <= page <= total_pages:
                start_idx = (page - 1) * paragraphs_per_page
                end_idx = min(start_idx + paragraphs_per_page, len(all_paragraphs))
                
                page_content = "\n".join(all_paragraphs[start_idx:end_idx])
                
                # 添加表格内容（如果有）
                if all_tables:
                    page_content += "\n\n--- 表格内容 ---\n"
                    for i, table in enumerate(all_tables):
                        page_content += f"\n表格 {i+1}:\n"
                        for row in table:
                            page_content += " | ".join(row) + "\n"
                
                return {
                    "preview_type": "word",
                    "total_pages": total_pages,
                    "current_page": page,
                    "content": page_content,
                    "message": f"第{page}页内容"
                }
            
            # 提取多页内容
            for i in range(min(max_pages, total_pages)):
                start_idx = i * paragraphs_per_page
                end_idx = min(start_idx + paragraphs_per_page, len(all_paragraphs))
                
                page_content = "\n".join(all_paragraphs[start_idx:end_idx])
                
                # 添加表格内容（如果有）
                if all_tables and i == 0:  # 只在第一页显示表格
                    page_content += "\n\n--- 表格内容 ---\n"
                    for j, table in enumerate(all_tables):
                        page_content += f"\n表格 {j+1}:\n"
                        for row in table:
                            page_content += " | ".join(row) + "\n"
                
                pages.append({
                    "page_number": i + 1,
                    "content": page_content
                })
            
            return {
                "preview_type": "word",
                "total_pages": total_pages,
                "pages_extracted": min(max_pages, total_pages),
                "pages": pages,
                "tables_count": len(all_tables),
                "message": f"成功提取{min(max_pages, total_pages)}页内容"
            }
        except Exception as e:
            return {
                "preview_type": "word",
                "error": f"处理Word文档时出错: {str(e)}",
                "pages": []
            }
    
    async def _preview_excel_document(
                                        self, 
                                        file_path: str, 
                                        max_sheets: int, 
                                        start_index: int = 0,
                                        sheet_index: int | None = None
                                    ) -> dict[str, Any]:
        """预览Excel文档 - 支持多种提取方式"""
        workbook = None
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            total_sheets = len(workbook.sheetnames)
            
            # 如果指定了单个工作表索引
            if sheet_index is not None:
                # 检查索引是否有效
                if sheet_index < 0 or sheet_index >= total_sheets:
                    return {
                        "preview_type": "excel",
                        "error": f"工作表索引 {sheet_index} 超出范围，总共有 {total_sheets} 个工作表",
                        "total_sheets": total_sheets
                    }
                
                return self._extract_single_sheet(workbook, sheet_index, total_sheets)
            
            # 提取多个工作表内容 - 从start_index开始提取
            sheets = []
            actual_start = max(0, min(start_index, total_sheets - 1))
            sheets_to_extract = min(max_sheets, total_sheets - actual_start)
            
            for i in range(actual_start, actual_start + sheets_to_extract):
                sheet_info = self._extract_sheet_preview(workbook, i, preview_rows=10)
                sheets.append(sheet_info)
            
            workbook.close()
            
            return {
                "preview_type": "excel",
                "content": sheets,
                "total_sheets": total_sheets,
                "start_index": actual_start,
                "sheets_extracted": sheets_to_extract,
                "message": f"成功提取从第{actual_start + 1}个开始共{sheets_to_extract}个工作表内容"
            }
            
        except Exception as e:
            if workbook:
                workbook.close()
            return {
                "preview_type": "excel",
                "error": f"处理Excel文档时出错: {str(e)}"
            }

    def _extract_single_sheet(self, workbook, sheet_index: int, total_sheets: int) -> dict:
        """提取单个工作表内容"""
        sheet_name = workbook.sheetnames[sheet_index]
        sheet = workbook[sheet_name]
        sheet_data = []
        
        # 读取前20行数据
        for row in sheet.iter_rows(max_row=20, values_only=True):
            if any(cell is not None for cell in row):
                sheet_data.append([str(cell) if cell is not None else "" for cell in row])
        
        workbook.close()
        
        return {
            "preview_type": "excel",
            "total_sheets": total_sheets,
            "current_sheet": sheet_name,
            "sheet_index": sheet_index,
            "data": sheet_data,
            "rows_extracted": len(sheet_data),
            "total_rows": sheet.max_row,
            "total_columns": sheet.max_column,
            "message": f"第{sheet_index + 1}个工作表 '{sheet_name}' 的内容"
        }

    def _extract_sheet_preview(self, workbook, sheet_index: int, preview_rows: int = 10) -> dict:
        """提取工作表预览信息"""
        sheet_name = workbook.sheetnames[sheet_index]
        sheet = workbook[sheet_name]
        sheet_data = []
        
        # 读取指定行数数据作为预览
        for row in sheet.iter_rows(max_row=preview_rows, values_only=True):
            if any(cell is not None for cell in row):
                sheet_data.append([str(cell) if cell is not None else "" for cell in row])
        
        return {
            "sheet_name": sheet_name,
            "sheet_index": sheet_index,
            "data": sheet_data,
            "rows_extracted": len(sheet_data),
            "total_rows": sheet.max_row,
            "total_columns": sheet.max_column
        }
    


    async def _preview_ppt_document(self, file_path: str, max_slides: int, slide_number: int | None) -> dict[str, Any]:
        """预览PPT文档 - 提取完整的幻灯片文件(包括相关资源)"""
        try:
            presentation = Presentation(file_path)
            total_slides = len(presentation.slides)
            
            # 创建临时目录来处理文件
            with tempfile.TemporaryDirectory() as temp_dir:
                # 如果指定了幻灯片编号，只提取该幻灯片
                if slide_number is not None and 1 <= slide_number <= total_slides:
                    # 创建只包含指定幻灯片的新PPTX
                    new_pptx_path = os.path.join(temp_dir, f"slide_{slide_number}.pptx")
                    
                    # 创建新演示文稿
                    new_presentation = Presentation()
                    
                    # 复制原幻灯片到新演示文稿
                    source_slide = presentation.slides[slide_number-1]
                    slide_layout = new_presentation.slide_layouts[6]  # 空白布局
                    new_slide = new_presentation.slides.add_slide(slide_layout)
                    
                    # 复制所有形状
                    for shape in source_slide.shapes:
                        if shape.has_text_frame:
                            # 复制文本框
                            new_shape = new_slide.shapes.add_textbox(
                                shape.left, shape.top, shape.width, shape.height
                            )
                            new_shape.text = shape.text # type: ignore
                        elif shape.has_table:
                            # 复制表格
                            table = shape.table # type: ignore
                            rows, cols = len(table.rows), len(table.columns)
                            new_table = new_slide.shapes.add_table(
                                rows, cols, shape.left, shape.top, shape.width, shape.height
                            ).table
                            # 复制表格内容
                            for i in range(rows):
                                for j in range(cols):
                                    if i < len(new_table.rows) and j < len(new_table.columns):
                                        new_table.cell(i, j).text = table.cell(i, j).text
                        # 可以添加更多形状类型的处理
                    
                    # 保存新演示文稿
                    new_presentation.save(new_pptx_path)
                    
                    # 读取文件内容并转换为base64
                    with open(new_pptx_path, 'rb') as f:
                        file_content = f.read()
                    
                    slide_base64 = base64.b64encode(file_content).decode('utf-8')
                    
                    return {
                        "preview_type": "ppt",
                        "total_slides": total_slides,
                        "current_slide": slide_number,
                        "content": slide_base64,
                        "message": f"第{slide_number}张幻灯片完整内容"
                    }
                
                # 提取多张幻灯片内容
                slides_content = []
                
                for i in range(min(max_slides, total_slides)):
                    # 为每个幻灯片创建单独的演示文稿
                    slide_pptx_path = os.path.join(temp_dir, f"slide_{i+1}.pptx")
                    
                    new_presentation = Presentation()
                    source_slide = presentation.slides[i]
                    slide_layout = new_presentation.slide_layouts[6]
                    new_slide = new_presentation.slides.add_slide(slide_layout)
                    
                    # 复制所有形状
                    for shape in source_slide.shapes:
                        if shape.has_text_frame:
                            new_shape = new_slide.shapes.add_textbox(
                                shape.left, shape.top, shape.width, shape.height
                            )
                            new_shape.text = shape.text # type: ignore
                        elif shape.has_table:
                            table = shape.table # type: ignore
                            rows, cols = len(table.rows), len(table.columns)
                            new_table = new_slide.shapes.add_table(
                                rows, cols, shape.left, shape.top, shape.width, shape.height
                            ).table
                            for row_idx in range(rows):
                                for col_idx in range(cols):
                                    if row_idx < len(new_table.rows) and col_idx < len(new_table.columns):
                                        new_table.cell(row_idx, col_idx).text = table.cell(row_idx, col_idx).text
                    
                    new_presentation.save(slide_pptx_path)
                    
                    # 读取文件内容并转换为base64
                    with open(slide_pptx_path, 'rb') as f:
                        file_content = f.read()
                    
                    slide_base64 = base64.b64encode(file_content).decode('utf-8')
                    
                    slides_content.append({
                        "slide_number": i + 1,
                        "content": slide_base64,
                        "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    })
                
                return {
                    "preview_type": "ppt",
                    "slides": slides_content,
                    "total_slides": total_slides,
                    "slides_extracted": min(max_slides, total_slides),
                    "message": f"成功提取{min(max_slides, total_slides)}张幻灯片完整内容"
                }
                
        except Exception as e:
            return {
                "preview_type": "ppt",
                "error": f"处理PPT文档时出错: {str(e)}",
                "slides": []
            }
    
    async def _preview_video(self, file_path: str) -> dict[str, Any]:
        """预览视频文件"""
        # 生成视频预览信息
        return {
            "preview_type": "video",
            "content": f"视频文件: {os.path.basename(file_path)}",
            "message": "视频文件需要在HTML5 video标签中播放"
        }
    
    async def _preview_audio(self, file_path: str) -> dict[str, Any]:
        """预览音频文件"""
        # 生成音频预览信息
        return {
            "preview_type": "audio",
            "content": f"音频文件: {os.path.basename(file_path)}",
            "message": "音频文件需要在HTML5 audio标签中播放"
        }
    